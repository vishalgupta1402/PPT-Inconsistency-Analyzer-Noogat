import os
import google.generativeai as genai
from pptx import Presentation
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Configure the Gemini API with your key
try:
    genai.configure(api_key=os.environ["GEMINI_API_KEY"])
except KeyError:
    print("Error: GEMINI_API_KEY not found. Please add your API key to the .env file.")
    exit()

def get_slide_content(file_path):
    """
    Opens a PowerPoint file and extracts text and table data from each slide.
    Returns a list of dictionaries, with each dictionary representing a slide.
    """
    try:
        prs = Presentation(file_path)
    except FileNotFoundError:
        print(f"Error: The file {file_path} was not found.")
        return None

    presentation_content = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": slide_number,
            "text_content": [],
            "table_content": []
        }

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join([run.text for run in paragraph.runs]).strip()
                    if text:
                        slide_data["text_content"].append(text)

            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_data["table_content"].append(table_data)

        presentation_content.append(slide_data)
        
    return presentation_content

def find_inconsistencies_with_gemini(all_slide_data):
    """
    Sends the presentation content to the Gemini API to find inconsistencies.
    """
    # Create the model
    model = genai.GenerativeModel('gemini-2.5-flash-preview-05-20')

    # Construct the prompt
    prompt_template = """
    You are an AI assistant specialized in finding factual and logical inconsistencies in presentation slide content.
    You will be provided with a JSON array of slide content. Each object in the array represents a slide and contains its slide number, text content, and table data.

    Your task is to analyze this content across all slides to identify any:
    1. Conflicting numerical data (e.g., revenue figures, percentages, time savings).
    2. Contradictory textual claims (e.g., "market is highly competitive" vs. "few competitors").
    3. Timeline mismatches (e.g., conflicting dates or forecasts).
    4. Any other logical inconsistencies.

    Please provide a clear and structured output in Markdown format. For each inconsistency you find, specify the following:
    - **Issue:** A brief description of the inconsistency.
    - **Slides:** The slide numbers involved.
    - **Details:** The specific conflicting information and where it is located.

    If you find no inconsistencies, state "No inconsistencies found."

    Here is the presentation content:
    {presentation_content}
    """
    
    # Format the prompt with the extracted data
    prompt = prompt_template.format(presentation_content=all_slide_data)
    
    print("\nSending data to Gemini API for analysis...")
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"An error occurred with the Gemini API: {e}"

if __name__ == "__main__":
    # Define file path relative to the script
    file_path = os.path.join(os.path.dirname(__file__), '..', 'presentations', 'NoogatAssignment.pptx')
    
    # Extract content from the presentation
    all_slide_data = get_slide_content(file_path)
    
    if all_slide_data:
        # Find inconsistencies using the Gemini API
        inconsistency_report = find_inconsistencies_with_gemini(all_slide_data)
        
        print("\n" + "="*50)
        print("INCONSISTENCY ANALYSIS REPORT")
        print("="*50)
        print(inconsistency_report)